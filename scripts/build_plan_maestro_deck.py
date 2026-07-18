from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "PLAN_MAESTRO_deck.pptx"

NAVY = RGBColor(16, 31, 54)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(93, 105, 122)
TEAL = RGBColor(18, 147, 133)
BLUE = RGBColor(55, 112, 180)
GOLD = RGBColor(235, 166, 54)
RED = RGBColor(202, 73, 73)
PALE = RGBColor(242, 246, 248)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(65, 145, 101)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def text(slide, s, x, y, w, h, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, font="DejaVu Sans", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = s; p.alignment = align
    r = p.runs[0]; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def rich(slide, runs, x, y, w, h, size=18, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for val, bold, col in runs:
        r = p.add_run(); r.text = val; r.font.name = "DejaVu Sans"; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = col or color
    return box


def base(title, kicker=None, num=None, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = NAVY if dark else WHITE
    if kicker:
        text(slide, kicker.upper(), .65, .35, 5.5, .3, 10, TEAL if not dark else GOLD, True)
    text(slide, title, .65, .72 if kicker else .42, 12, .7, 28, WHITE if dark else NAVY, True)
    if not dark: rect(slide, .65, 1.42, 1.0, .05, TEAL)
    if num is not None:
        text(slide, f"{num:02d}", 12.25, 7.05, .45, .2, 9, RGBColor(150,160,171), True, PP_ALIGN.RIGHT)
    return slide


def bullet_list(slide, items, x, y, w, h, size=17, color=INK, accent=TEAL, gap=0.52):
    for i, item in enumerate(items):
        cy = y + i * gap
        rect(slide, x, cy+.12, .09, .09, accent, True)
        text(slide, item, x+.22, cy, w-.22, gap+.05, size, color)


def card(slide, title_, body, x, y, w, h, accent=TEAL, stat=None):
    rect(slide, x, y, w, h, PALE, True, RGBColor(222,229,233))
    rect(slide, x, y, .08, h, accent, True)
    text(slide, title_, x+.28, y+.2, w-.5, .35, 15, NAVY, True)
    if stat:
        text(slide, stat, x+.28, y+.65, w-.5, .55, 28, accent, True)
        text(slide, body, x+.28, y+1.25, w-.5, h-1.4, 12, MUTED)
    else:
        text(slide, body, x+.28, y+.68, w-.5, h-.85, 12.5, MUTED)


# 1 — cover
s = base("Dinero público → resultados", dark=True)
text(s, "PLAN MAESTRO", .72, .58, 5, .35, 12, GOLD, True)
text(s, "Un marco para medir evolución, rendimiento y escenarios fiscales", .72, 2.05, 8.8, 1.25, 30, WHITE, True)
text(s, "TFM · versión ejecutiva", .75, 3.55, 4.5, .35, 16, RGBColor(191,205,219))
rect(s, 9.7, 1.1, 2.45, 4.8, TEAL, True)
text(s, "1900", 10.15, 1.55, 1.6, .5, 25, WHITE, True, PP_ALIGN.CENTER)
text(s, "↓", 10.6, 2.35, .7, .6, 30, GOLD, True, PP_ALIGN.CENTER)
text(s, "2023", 10.15, 3.18, 1.6, .5, 25, WHITE, True, PP_ALIGN.CENTER)
text(s, "datos · modelos\nincertidumbre", 10.02, 4.35, 1.85, .8, 14, WHITE, False, PP_ALIGN.CENTER)
text(s, "18 julio 2026", .75, 6.8, 3, .3, 11, RGBColor(160,178,194))

# 2 — thesis
s = base("La pregunta es simple. La respuesta no puede ser una liga.", "tesis", 2)
rich(s, [("¿Qué país convierte mejor el dinero público en resultados", True, NAVY),
         (" — dadas sus condiciones y la incertidumbre?", False, MUTED)], .8, 1.85, 11.7, .8, 25)
card(s, "Medir", "Resultados observados frente a resultados esperados, controlando renta, demografía, gobernanza y capacidad estadística.", .8, 3.05, 3.6, 2.25, BLUE)
card(s, "Comparar", "Funnel plots por grupo de renta e intervalos conformales: posiciones, no podios.", 4.85, 3.05, 3.6, 2.25, TEAL)
card(s, "Proyectar", "Escenarios condicionados con consecuencias y deuda; nunca una recomendación política única.", 8.9, 3.05, 3.6, 2.25, GOLD)
text(s, "Principio rector: rendimiento ajustado ≠ causalidad ≠ eficiencia absoluta", .9, 6.3, 11.5, .4, 17, RED, True, PP_ALIGN.CENTER)

# 3 scope
s = base("Un programa amplio; un TFM deliberadamente acotado", "decisión de alcance", 3)
card(s, "Núcleo demostrable", "Vista C vivienda o módulo global de sanidad. Una pregunta, un outcome emparejado, un funnel defendible.", .75, 1.75, 3.8, 3.65, TEAL)
card(s, "Capítulo descriptivo", "Atlas de evolución: 15 preguntas, figuras comparables y caveats explícitos. Sin forzar series históricas inexistentes.", 4.77, 1.75, 3.8, 3.65, BLUE)
card(s, "Predicción", "Panel trimestral CCAA, baselines obligatorios y backtesting. DL solo como comparador.", 8.79, 1.75, 3.8, 3.65, GOLD)
rect(s, .75, 5.8, 11.84, .76, NAVY, True)
text(s, "F5–F9 quedan como continuación natural del programa, salvo requisito expreso de la vista elegida.", 1.0, 6.0, 11.3, .3, 16, WHITE, True, PP_ALIGN.CENTER)

# 4 research architecture
s = base("Cuatro bloques conectan descripción, explicación y decisión", "arquitectura", 4)
labels = [("B", "EVOLUCIÓN", "1900–2023\n15 series"), ("A", "RENDIMIENTO", "gasto → outcome\najustado"),
          ("C", "PREDICCIÓN", "forecast +\nescenarios"), ("D", "POLÍTICA", "menú de\nconsecuencias")]
xs = [.65, 3.75, 6.85, 9.95]
cols = [BLUE, TEAL, GOLD, RED]
for i, (letter, label, sub) in enumerate(labels):
    rect(s, xs[i], 2.0, 2.7, 2.65, PALE, True, cols[i])
    text(s, letter, xs[i]+.15, 2.25, .6, .6, 30, cols[i], True, PP_ALIGN.CENTER)
    text(s, label, xs[i]+.75, 2.35, 1.75, .35, 14, NAVY, True)
    text(s, sub, xs[i]+.35, 3.25, 2.0, .75, 14, MUTED, False, PP_ALIGN.CENTER)
    if i < 3: text(s, "→", xs[i]+2.72, 2.92, .38, .4, 20, MUTED, True, PP_ALIGN.CENTER)
text(s, "La elección política permanece fuera del modelo", 3.2, 5.45, 6.95, .5, 19, RED, True, PP_ALIGN.CENTER)
text(s, "El trabajo estima patrones, incertidumbre y trade-offs; no decide qué debe despriorizarse.", 2.2, 6.0, 8.95, .45, 14, MUTED, False, PP_ALIGN.CENTER)

# 5 data
s = base("La base de datos ya existe y ha pasado controles", "preparación", 5)
stats = [("28", "datasets processed", TEAL), ("4", "tablas gold", BLUE), ("202", "países en GMD", GOLD), ("1900–2030", "historia + ancla WEO", RED)]
for i,(st,lab,col) in enumerate(stats):
    x=.7+i*3.1; rect(s,x,1.75,2.75,1.35,PALE,True); text(s,st,x+.15,1.95,2.45,.5,26,col,True,PP_ALIGN.CENTER); text(s,lab,x+.15,2.55,2.45,.25,11,MUTED,False,PP_ALIGN.CENTER)
bullet_list(s,["UE: gasto funcional × económico, ingresos, déficit y deuda (1995–2023)",
               "Global: WEO, GHED, outcomes sociales y demográficos",
               "Histórico: GMD + JST; 158 países con observaciones pre-1995",
               "España: panel trimestral, CCAA, precios, salarios y vivienda"],.95,3.65,7.2,2.3,15,gap=.55)
card(s,"Pendiente inmediato","Desempleo · intereses D41 · pensiones y población ≥65 · extras históricos D62",8.65,3.55,3.75,2.15,GOLD)
text(s,"Regla de calidad: nada entra en gold sin smoke test",8.8,6.05,3.45,.4,14,TEAL,True,PP_ALIGN.CENTER)

# 6 atlas
s = base("El atlas cuenta “el siglo del Estado” con una gramática común", "bloque B", 6)
groups = [("Tamaño fiscal", "gasto · ingresos · déficit · deuda · intereses"),
          ("Composición", "salarios · capital · vivienda · sanidad · protección"),
          ("Presiones", "paro · pensiones · migración · ayuda internacional")]
for i,(a,b) in enumerate(groups): card(s,a,b,.75+i*4.05,1.75,3.7,1.55,[BLUE,TEAL,GOLD][i])
text(s,"Plantilla de cada figura",.8,3.75,3.2,.4,17,NAVY,True)
bullet_list(s,["mediana mundial", "selección estable de países", "España destacada", "eras y rupturas anotadas"],.9,4.25,4.2,2.2,15,gap=.43)
rect(s,5.35,3.72,6.85,2.3,PALE,True)
text(s,"10,2%",5.65,4.07,2.0,.55,27,BLUE,True,PP_ALIGN.CENTER)
text(s,"→",7.78,4.08,.6,.5,24,MUTED,True,PP_ALIGN.CENTER)
text(s,"29,8%",8.48,4.07,2.0,.55,27,TEAL,True,PP_ALIGN.CENTER)
text(s,"mediana mundial del gasto público: primera era → era reciente",5.75,4.9,5.95,.45,14,MUTED,False,PP_ALIGN.CENTER)
text(s,"No se interpola un desglose funcional continuo antes de ~1970.",5.75,5.5,5.95,.35,13,RED,True,PP_ALIGN.CENTER)

# 7 method pipeline
s = base("El rendimiento ajustado se estima dentro de un protocolo auditable", "método", 7)
steps=[("1","Preparar","medias 5a · retardos 3–5a"),("2","Predecir","OLS/GAM → GBM"),("3","Ajustar","residual = observado − esperado"),("4","Calibrar","conformal por renta"),("5","Auditar","LOCO · bloques · multiverso")]
for i,(n,a,b) in enumerate(steps):
    x=.55+i*2.55; rect(s,x,2.0,2.25,2.4,PALE,True); text(s,n,x+.78,2.2,.7,.55,26,[BLUE,TEAL,GOLD,RED,NAVY][i],True,PP_ALIGN.CENTER); text(s,a,x+.2,3.02,1.85,.32,14,NAVY,True,PP_ALIGN.CENTER); text(s,b,x+.18,3.48,1.9,.6,11,MUTED,False,PP_ALIGN.CENTER)
    if i<4:text(s,"→",x+2.25,2.92,.3,.35,17,MUTED,True,PP_ALIGN.CENTER)
rect(s,.75,5.05,11.85,.95,NAVY,True)
text(s,"Criterio de aceptación: residual ⟂ renta y residual ⟂ SPI dentro de cada grupo",1.0,5.3,11.35,.35,17,WHITE,True,PP_ALIGN.CENTER)
text(s,"Encuadre no causal · intervalos visibles · estabilidad Spearman reportada",2.1,6.38,9.15,.35,14,MUTED,False,PP_ALIGN.CENTER)

# 8 ML map
s = base("El tamaño efectivo de muestra decide la herramienta", "ML / DL", 8)
rows=[("GBM + SHAP + conformal","rendimiento y composición","N panel suficiente",TEAL),
      ("PCA · clustering · DTW","tipologías y trayectorias","siglo fiscal",BLUE),
      ("SARIMAX / GBM","forecast trimestral CCAA","rolling backtest",GOLD),
      ("N-BEATS / DeepAR","comparador, no protagonista","solo si bate baseline",RED)]
for i,(a,b,c,col) in enumerate(rows):
    y=1.75+i*1.02; rect(s,.72,y,11.9,.76,PALE,True); rect(s,.72,y,.1,.76,col); text(s,a,1.05,y+.19,3.2,.3,14,NAVY,True); text(s,b,4.45,y+.19,4.1,.3,13,INK); text(s,c,9.0,y+.19,3.1,.3,12,MUTED)
text(s,"Descartes explícitos",.8,6.0,2.3,.3,14,RED,True)
text(s,"sin LSTM/TFT primario · sin GNN a ~150 nodos · sin LLM como extractor masivo · sin liga mundial",2.75,5.96,9.5,.5,13,MUTED)

# 9 prediction
s = base("Predicción: competir contra lo simple y publicar el resultado", "bloque C", 9)
card(s,"Modo 1 · Forecast real","Panel trimestral CCAA. Naive/ETS → SARIMAX/GBM → N-BEATS/DeepAR. Evaluación rolling-origin a 4–8 trimestres.",.75,1.7,3.7,3.75,TEAL)
card(s,"Modo 2 · Proyección condicionada","Sendas alternativas de gasto y composición. El modelo A1 proyecta outcomes con intervalos conformales.",4.8,1.7,3.7,3.75,BLUE)
card(s,"Modo 3 · Ancla externa","Proyecciones FMI-WEO hasta 2030 para contraste y consistencia macro.",8.85,1.7,3.7,3.75,GOLD)
text(s,"“Batir al naive es noticia; no batirlo también se publica.”",1.3,6.15,10.75,.45,20,NAVY,True,PP_ALIGN.CENTER)

# 10 scenarios
s = base("Del “plan fiscal” al simulador de escenarios", "bloque D", 10)
text(s,"ENTRADAS",.9,1.8,2.4,.35,13,MUTED,True,PP_ALIGN.CENTER)
text(s,"MOTOR",5.45,1.8,2.4,.35,13,MUTED,True,PP_ALIGN.CENTER)
text(s,"SALIDAS",10.0,1.8,2.4,.35,13,MUTED,True,PP_ALIGN.CENTER)
card(s,"Sendas alternativas","+0,5 pp vivienda\n+0,5 pp inversión\n−X pp partida Y\nsaldo primario Z",.75,2.3,3.0,2.8,BLUE)
text(s,"→",3.85,3.42,.55,.55,27,MUTED,True,PP_ALIGN.CENTER)
card(s,"Dos modelos","Aritmética de deuda r−g\n+\nProyección condicionada de resultados",4.5,2.3,4.0,2.8,TEAL)
text(s,"→",8.62,3.42,.55,.55,27,MUTED,True,PP_ALIGN.CENTER)
card(s,"Matriz de consecuencias","deuda 2024–2035\noutcomes esperados\nintervalos\ntrade-offs",9.28,2.3,3.0,2.8,GOLD)
rect(s,1.45,5.75,10.45,.72,NAVY,True)
text(s,"Producto científico: menú comparable. Elección final: política.",1.75,5.95,9.85,.3,17,WHITE,True,PP_ALIGN.CENTER)

# 11 roadmap
s = base("Diez semanas, tres gates y un buffer real", "plan de ejecución", 11)
weeks=[("S1","gate tutor",RED),("S2–3","entregas + atlas",BLUE),("S4–5","sanidad + funnel",TEAL),("S6","vivienda + cierre B",GOLD),("S7","forecast + síntesis",TEAL),("S8","app + memoria",BLUE),("S9","borrador",RED),("S10","buffer",NAVY)]
x=.6
for lab,desc,col in weeks:
    w=1.35 if len(lab)<3 else 1.6
    rect(s,x,2.05,w,1.55,col,True)
    text(s,lab,x+.08,2.3,w-.16,.35,16,WHITE,True,PP_ALIGN.CENTER)
    text(s,desc,x+.08,2.82,w-.16,.5,10.5,WHITE,False,PP_ALIGN.CENTER)
    x+=w+.12
text(s,"Gates",.8,4.55,1.2,.3,14,NAVY,True)
bullet_list(s,["S1: aval escrito y vista elegida", "S5: primer funnel y checkpoint", "S9: borrador completo para revisión"],.9,5.05,5.8,1.5,14,gap=.43)
card(s,"Línea roja","Si el alcance vuelve a crecer, se sacrifica profundidad, reproducibilidad y defensa.",7.55,4.68,4.75,1.62,RED)

# 12 deliverables
s = base("Cinco productos convierten el análisis en una entrega coherente", "salidas", 12)
items=[("01","Memoria TFM","método · límites · reproducibilidad"),("02","Atlas","15 figuras + tabla de eras"),("03","Funnels","rendimiento ajustado por módulo"),("04","Simulador","deuda + composición + outcomes"),("05","App + gold","exploración reproducible")]
for i,(n,a,b) in enumerate(items):
    y=1.65+i*.98; text(s,n,.75,y,.6,.35,14,TEAL,True); text(s,a,1.45,y,2.2,.35,16,NAVY,True); text(s,b,3.85,y,7.7,.35,14,MUTED); rect(s,.75,y+.55,11.55,.02,RGBColor(225,230,234))

# 13 risks
s = base("Los riesgos están en primera plana — con controles concretos", "defensa", 13)
risks=[("Normatividad","simular; no prescribir",RED),("Endogeneidad","retardos + encuadre no causal",GOLD),("Sesgo del panel","procedencia + specs sin imputados",BLUE),("Comparabilidad","histórico solo descriptivo",TEAL),("Goodhart","intervalos + no-liga",RED),("P-hacking","RQs y métricas pre-registradas",NAVY)]
for i,(a,b,col) in enumerate(risks):
    colx=.7+(i%2)*6.05; y=1.65+(i//2)*1.62
    rect(s,colx,y,5.65,1.22,PALE,True); rect(s,colx,y,.09,1.22,col,True)
    text(s,a,colx+.32,y+.18,1.75,.32,14,NAVY,True); text(s,b,colx+2.1,y+.18,3.2,.7,13,MUTED)
text(s,"Límite que no desaparece: el residual también absorbe instituciones, historia, geografía e informalidad.",1.1,6.68,11.1,.32,13,RED,True,PP_ALIGN.CENTER)

# 14 decision
s = base("La decisión inmediata es de alcance, no de algoritmo", "siguiente paso", 14)
rect(s,.75,1.7,11.85,1.0,NAVY,True)
text(s,"Obtener el aval escrito del tutor y fijar una única vista central",1.0,2.0,11.35,.35,20,WHITE,True,PP_ALIGN.CENTER)
card(s,"Opción recomendada","Vivienda UE como núcleo complementario del marco integral.",.9,3.35,3.55,1.7,TEAL)
card(s,"Alternativa sólida","Sanidad global: mayor N y módulo A1 más directo.",4.9,3.35,3.55,1.7,BLUE)
card(s,"Después del gate","Re-pulls menores → atlas núcleo → primer funnel.",8.9,3.35,3.55,1.7,GOLD)
text(s,"Criterio de éxito del TFM: una afirmación modesta, reproducible y difícil de desmontar.",1.25,6.15,10.85,.45,18,NAVY,True,PP_ALIGN.CENTER)

# 15 close
s = base("Relevante y original. Defendible si se mantiene el recorte.", dark=True)
text(s,"Tres compromisos",.78,1.7,3.2,.35,13,GOLD,True)
for i,t in enumerate(["alcance TFM radicalmente acotado","escenarios, no prescripción","causalidad limitada en primera plana"]):
    rect(s,.82,2.35+i*1.05,.32,.32,TEAL,True)
    text(s,t,1.35,2.25+i*1.05,7.9,.55,21,WHITE,True)
text(s,"Dinero público → resultados",.82,6.45,6.5,.35,14,RGBColor(169,188,204))
rect(s,10.35,1.55,1.7,4.2,TEAL,True)
text(s,"MEDIR",10.55,2.05,1.3,.35,15,WHITE,True,PP_ALIGN.CENTER)
text(s,"↓",10.85,2.85,.7,.45,24,GOLD,True,PP_ALIGN.CENTER)
text(s,"DUDAR",10.55,3.65,1.3,.35,15,WHITE,True,PP_ALIGN.CENTER)
text(s,"↓",10.85,4.35,.7,.45,24,GOLD,True,PP_ALIGN.CENTER)
text(s,"DECIDIR",10.45,5.02,1.5,.35,15,WHITE,True,PP_ALIGN.CENTER)


# Metadata and save
prs.core_properties.title = "Dinero público → resultados — Plan Maestro"
prs.core_properties.subject = "Deck ejecutivo basado en PLAN_MAESTRO.md"
prs.core_properties.author = "Proyecto evo_final_work"
prs.core_properties.comments = "Generado desde docs/PLAN_MAESTRO.md el 2026-07-18"
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
